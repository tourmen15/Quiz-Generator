import os
import base64
import io
import random
import re
import uuid
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from transformers import pipeline
import PyPDF2
import docx
import pptx
from fpdf import FPDF

# --- Configuration & Initialization ---

app = Flask(__name__)
# Allow requests from the React frontend which will be running on a different port
CORS(app) 

# Initialize the Hugging Face pipeline for question generation.
# This is loaded once when the application starts to avoid reloading the model on every request,
# which would be very slow.
try:
    print("Loading Hugging Face model... This may take a few minutes on first run.")
    question_generator = pipeline("text2text-generation", model="iarfmoose/t5-base-question-generator")
    print("Model loaded successfully.")
except Exception as e:
    # If the model fails to load, we can't generate questions. 
    # Log the error and handle it gracefully in the endpoint.
    print(f"Error loading Hugging Face model: {e}")
    question_generator = None

# --- Document Parsing Utilities ---

def parse_pdf(file_content):
    """Extracts text from a PDF file's content."""
    text = ""
    try:
        pdf_file = io.BytesIO(file_content)
        reader = PyPDF2.PdfReader(pdf_file)
        # Enforce a strict 200-page limit as per requirements
        if len(reader.pages) > 200:
            raise ValueError("File exceeds the 200-page limit.")
        for page in reader.pages:
            text += page.extract_text() or ""
    except Exception as e:
        # Catch potential errors from corrupted PDFs
        raise ValueError(f"Error parsing PDF file: {e}")
    return text

def parse_docx(file_content):
    """Extracts text from a DOCX file's content."""
    text = ""
    try:
        doc_file = io.BytesIO(file_content)
        document = docx.Document(doc_file)
        for para in document.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        raise ValueError(f"Error parsing DOCX file: {e}")
    return text

def parse_pptx(file_content):
    """Extracts text from a PPTX file's content."""
    text = ""
    try:
        ppt_file = io.BytesIO(file_content)
        presentation = pptx.Presentation(ppt_file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        raise ValueError(f"Error parsing PPTX file: {e}")
    return text

def parse_txt(file_content):
    """Extracts text from a TXT file's content."""
    try:
        return file_content.decode('utf-8')
    except Exception as e:
        raise ValueError(f"Error parsing TXT file: {e}")

# --- AI & Quiz Generation Logic ---

def get_text_chunks(full_text):
    """
    Splits the full text into smaller, more manageable paragraphs.
    This provides better, focused context for the AI model.
    """
    # Split by double newlines, then filter out empty or very short chunks
    paragraphs = full_text.split('\n\n')
    return [p.strip() for p in paragraphs if len(p.strip()) > 150]

def get_keywords(text_chunk):
    """
    A simple heuristic to find potential keywords (answers) in a text chunk.
    This method looks for capitalized words, which often represent named entities or key concepts.
    It's a simplified approach to avoid adding heavy dependencies like SpaCy or NLTK.
    """
    # Use regex to find capitalized words, but not those at the start of a sentence.
    # This is a bit tricky, so we'll settle for finding all capitalized words and filtering common ones.
    keywords = re.findall(r'\b[A-Z][a-z]+\b', text_chunk)
    # Filter out common English words that might be capitalized by mistake
    common_words = {"The", "A", "An", "Is", "Was", "Were"}
    return list(set([kw for kw in keywords if kw not in common_words]))

def create_quiz_instance(full_text, num_questions, question_types):
    """
    Generates a single version of a quiz with a specified number and type of questions.
    """
    if not question_generator:
        raise RuntimeError("Question generation model is not available.")

    chunks = get_text_chunks(full_text)
    if not chunks:
        raise ValueError("Could not extract any meaningful text chunks from the provided content.")

    questions = []
    used_chunks_indices = set() # To promote variety in questions

    while len(questions) < num_questions:
        if len(used_chunks_indices) >= len(chunks):
             # If we've used all chunks and still need questions, allow reuse.
             used_chunks_indices.clear()
        
        # Select a random chunk of text that hasn't been used yet, if possible
        available_indices = [i for i in range(len(chunks)) if i not in used_chunks_indices]
        if not available_indices:
            if not chunks: break # No text chunks to process
            # Fallback to any chunk if all have been used
            chunk_index = random.randint(0, len(chunks) - 1)
        else:
             chunk_index = random.choice(available_indices)

        used_chunks_indices.add(chunk_index)
        context = chunks[chunk_index]
        
        keywords = get_keywords(context)
        if not keywords:
            continue # Skip chunks where we can't find good keywords

        correct_answer = random.choice(keywords)

        # Decide question type for this iteration
        current_question_type = question_types
        if question_types == 'mixed':
            current_question_type = random.choice(['mcq_only', 'fill_only'])

        try:
            if current_question_type == 'mcq_only':
                # Generate the question using the AI model
                input_text = f"answer: {correct_answer} context: {context}"
                generated = question_generator(input_text, max_length=128)
                question_text = generated[0]['generated_text'].replace("question: ", "").strip()

                # Generate distractors (plausible but incorrect answers)
                all_keywords = list(set([kw for chunk in chunks for kw in get_keywords(chunk)]))
                distractors = [kw for kw in all_keywords if kw != correct_answer]
                
                # Ensure we have enough distractors
                if len(distractors) < 3:
                    distractors.extend(["Incorrect Option A", "Incorrect Option B", "Incorrect Option C"]) # Fallback

                random.shuffle(distractors)
                options_list = [correct_answer] + distractors[:3]
                random.shuffle(options_list)

                # Format options as A, B, C, D
                options_dict = {chr(65 + i): opt for i, opt in enumerate(options_list)}
                
                questions.append({
                    "id": str(uuid.uuid4()),
                    "type": "mcq",
                    "question_text": question_text,
                    "options": options_dict,
                    "correct_answer": correct_answer,
                })

            elif current_question_type == 'fill_only':
                # Find the sentence containing the answer and create a blank
                sentences = re.split(r'(?<=[.!?]) +', context)
                question_sentence = ""
                for sent in sentences:
                    if correct_answer in sent:
                        question_sentence = sent
                        break
                
                if not question_sentence: continue # Could not find the sentence

                question_text = question_sentence.replace(correct_answer, "___")
                
                questions.append({
                    "id": str(uuid.uuid4()),
                    "type": "fill_in_the_blank",
                    "question_text": question_text,
                    "correct_answer": correct_answer,
                })
        except Exception as e:
            print(f"Error generating a question: {e}")
            # Continue to the next attempt
            continue

    return questions

# --- File Exporting Utilities ---

class PDF(FPDF):
    def header(self):
        self.set_font('Arial', 'B', 12)
        self.cell(0, 10, 'Generated Quiz', 0, 1, 'C')

    def footer(self):
        self.set_y(-15)
        self.set_font('Arial', 'I', 8)
        self.cell(0, 10, f'Page {self.page_no()}', 0, 0, 'C')

def export_pdf(quiz_data):
    pdf = PDF()
    pdf.add_page()
    pdf.set_font('Arial', '', 11)
    
    for i, q in enumerate(quiz_data):
        # Using multi_cell for automatic line breaking
        pdf.multi_cell(0, 5, f"Q{i+1}: {q['question_text']}")
        pdf.ln(5)
        
        if q['type'] == 'mcq' and 'options' in q:
            for key, val in q['options'].items():
                pdf.multi_cell(0, 5, f"   {key}) {val}")
            pdf.ln(2)
            pdf.multi_cell(0, 5, f"   Correct Answer: {q['correct_answer']}")

        elif q['type'] == 'fill_in_the_blank':
            pdf.multi_cell(0, 5, f"   Correct Answer: {q['correct_answer']}")
        
        pdf.ln(10) # Space between questions

    pdf_bytes = pdf.output(dest='S').encode('latin-1')
    return io.BytesIO(pdf_bytes)

def export_docx(quiz_data):
    doc = docx.Document()
    doc.add_heading('Generated Quiz', 0)

    for i, q in enumerate(quiz_data):
        doc.add_paragraph(f"Q{i+1}: {q['question_text']}", style='List Number')
        
        if q['type'] == 'mcq' and 'options' in q:
            for key, val in q['options'].items():
                doc.add_paragraph(f"   {key}) {val}", style='List Bullet')
            doc.add_paragraph(f"   Correct Answer: {q['correct_answer']}")

        elif q['type'] == 'fill_in_the_blank':
             doc.add_paragraph(f"   Correct Answer: {q['correct_answer']}")
        
        doc.add_paragraph() # Space between questions

    doc_io = io.BytesIO()
    doc.save(doc_io)
    doc_io.seek(0)
    return doc_io

def export_txt(quiz_data):
    txt = ""
    for i, q in enumerate(quiz_data):
        txt += f"Q{i+1}: {q['question_text']}\n"
        
        if q['type'] == 'mcq' and 'options' in q:
            for key, val in q['options'].items():
                txt += f"   {key}) {val}\n"
            txt += f"   Correct Answer: {q['correct_answer']}\n"

        elif q['type'] == 'fill_in_the_blank':
            txt += f"   Correct Answer: {q['correct_answer']}\n"
        
        txt += "\n"

    return io.BytesIO(txt.encode('utf-8'))

# --- API Endpoints ---

@app.route('/api/generate-quiz', methods=['POST'])
def generate_quiz():
    """
    Main endpoint to handle quiz generation requests from the frontend.
    """
    try:
        data = request.get_json()
        source_type = data.get('source_type')
        content = data.get('content')
        num_questions = int(data.get('num_questions', 10))
        question_types = data.get('question_types', 'mcq_only')
        num_versions = int(data.get('num_versions', 1))

        if not content:
            return jsonify({"message": "Content is missing."}), 400

        full_text = ""
        if source_type == 'text':
            # Enforce 100,000 character limit for raw text
            if len(content) > 100000:
                return jsonify({"message": "Text input exceeds the 100,000 character limit."}), 400
            full_text = content
        elif source_type == 'file':
            file_name = data.get('file_name', '')
            file_content = base64.b64decode(content)
            
            ext = os.path.splitext(file_name)[1].lower()
            if ext == '.pdf':
                full_text = parse_pdf(file_content)
            elif ext == '.docx':
                full_text = parse_docx(file_content)
            elif ext == '.pptx':
                full_text = parse_pptx(file_content)
            elif ext == '.txt':
                full_text = parse_txt(file_content)
            else:
                return jsonify({"message": f"Unsupported file type: {ext}"}), 400
        else:
             return jsonify({"message": "Invalid source_type specified."}), 400

        # AI Generation Loop
        quizzes = []
        for i in range(num_versions):
            quiz_instance = create_quiz_instance(full_text, num_questions, question_types)
            quizzes.append({
                "version": i + 1,
                "questions": quiz_instance
            })

        return jsonify({"quizzes": quizzes}), 200

    except ValueError as ve:
        # Handle specific data validation and parsing errors
        return jsonify({"message": str(ve)}), 400
    except Exception as e:
        # Generic server error
        print(f"An unexpected error occurred: {e}")
        return jsonify({"message": "An internal server error occurred while generating the quiz."}), 500


@app.route('/api/export-quiz', methods=['POST'])
def export_quiz():
    """
    Endpoint to handle exporting the user-edited quiz to a file.
    """
    try:
        data = request.get_json()
        quiz_data = data.get('quiz_data')
        export_format = data.get('format')

        if not quiz_data or not export_format:
            return jsonify({"message": "Missing quiz data or format."}), 400

        file_io = None
        mimetype = ''
        filename = ''

        if export_format == 'pdf':
            file_io = export_pdf(quiz_data)
            mimetype = 'application/pdf'
            filename = 'quiz.pdf'
        elif export_format == 'docx':
            file_io = export_docx(quiz_data)
            mimetype = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            filename = 'quiz.docx'
        elif export_format == 'txt':
            file_io = export_txt(quiz_data)
            mimetype = 'text/plain'
            filename = 'quiz.txt'
        else:
            return jsonify({"message": "Unsupported export format."}), 400

        return send_file(
            file_io,
            mimetype=mimetype,
            as_attachment=_True,
            download_name=filename
        )

    except Exception as e:
        print(f"An unexpected error occurred during export: {e}")
        return jsonify({"message": "An internal server error occurred during file export."}), 500


if __name__ == '__main__':
    # Runs the Flask app. Debug=True allows for hot-reloading during development.
    app.run(debug=True, port=5000)
